import collections 
import collections.abc
from pptx import Presentation
from pptx import Presentation
import os
from PIL import Image
import re
def sorted_alphanumeric(data):
    convert = lambda text: int(text) if text.isdigit() else text.lower()
    alphanum_key = lambda key: [ convert(c) for c in re.split('([0-9]+)', key) ] 
    return sorted(data, key=alphanum_key)
def _add_image(slide, placeholder_id, image_url):
    placeholder = slide.placeholders[placeholder_id]

    # Calculate the image size of the image
    im = Image.open(image_url)
    width, height = im.size

    # Make sure the placeholder doesn't zoom in
    placeholder.height = height
    placeholder.width = width

    # Insert the picture
    placeholder = placeholder.insert_picture(image_url)

    # Calculate ratios and compare
    image_ratio = width / height
    placeholder_ratio = placeholder.width / placeholder.height
    ratio_difference = placeholder_ratio - image_ratio

    # Placeholder width too wide:
    if ratio_difference > 0:
        difference_on_each_side = ratio_difference / 5
        placeholder.crop_left = -difference_on_each_side
        placeholder.crop_right = -difference_on_each_side
    # Placeholder height too high
    else:
        difference_on_each_side = -ratio_difference / 5
        placeholder.crop_bottom = -difference_on_each_side
        placeholder.crop_top = -difference_on_each_side


prs = Presentation()
directory = 'Rapor Süreci' #folder name of images
lst = sorted_alphanumeric(os.listdir(directory))
for filename in lst:
    f = os.path.join(directory, filename)
    if os.path.isfile(f):
        print(f)
        layout8 = prs.slide_layouts[8]
        slide = prs.slides.add_slide(layout8)
        dosya = os.path.basename(f)
        title = slide.shapes.title.text = os.path.splitext(dosya)[0] 
        #sub = slide.placeholders[2].text = "Python has the power"
        _add_image(slide,1,f)

prs.save("MyPresentation.pptx")
os.startfile("MyPresentation.pptx")
